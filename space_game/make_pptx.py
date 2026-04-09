"""Generates Space Defender teaching slides (space_game_slides.pptx)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Colour palette (matches the game) ─────────────────────────────────────────
BG       = RGBColor(0x0A, 0x0A, 0x1E)   # dark navy
CYAN     = RGBColor(0x00, 0xFF, 0xFF)
YELLOW   = RGBColor(0xFF, 0xFF, 0x00)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE   = RGBColor(0xFF, 0x8C, 0x00)
RED      = RGBColor(0xDC, 0x28, 0x28)
GREEN    = RGBColor(0x00, 0xDC, 0x50)
PURPLE   = RGBColor(0xA0, 0x20, 0xF0)
GREY     = RGBColor(0xA0, 0xA0, 0xA0)
DK_GREY  = RGBColor(0x1A, 0x1A, 0x3A)  # slightly lighter than BG for code bg

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]   # completely blank

# ── Helpers ────────────────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(blank_layout)

def bg(slide):
    """Fill slide background with dark navy."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG

def textbox(slide, text, l, t, w, h,
            size=24, bold=False, color=WHITE,
            align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return txb

def code_box(slide, code, l, t, w, h, size=14):
    """Dark rounded rectangle + monospace code."""
    shape = slide.shapes.add_shape(
        1,   # MSO_SHAPE_TYPE.RECTANGLE
        l, t, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DK_GREY
    shape.line.color.rgb = CYAN
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    first = True
    for line in code.strip().split("\n"):
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line if line else " "
        run.font.name = "Courier New"
        run.font.size = Pt(size)
        run.font.color.rgb = CYAN

def divider(slide, y, color=CYAN, width_in=12):
    """Horizontal accent line."""
    ln = slide.shapes.add_shape(1,
        Inches(0.67), y,
        Inches(width_in), Pt(2))
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()

def bullet_list(slide, items, l, t, w, h, size=20, color=WHITE, indent_color=CYAN):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = True
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        # bullet dot
        dot = p.add_run()
        dot.text = "▸  "
        dot.font.size  = Pt(size)
        dot.font.color.rgb = indent_color
        dot.font.bold  = True
        # text
        run = p.add_run()
        run.text = item
        run.font.size  = Pt(size)
        run.font.color.rgb = color

def tag(slide, text, l, t, color=CYAN, bg_color=None):
    """Small coloured label/tag."""
    if bg_color:
        box = slide.shapes.add_shape(1, l, t, Inches(2.2), Inches(0.38))
        box.fill.solid()
        box.fill.fore_color.rgb = bg_color
        box.line.fill.background()
    txb = slide.shapes.add_textbox(l + Inches(0.1), t + Pt(3), Inches(2), Inches(0.35))
    tf  = txb.text_frame
    p   = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(13)
    run.font.bold  = True
    run.font.color.rgb = color

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)

# Decorative stars (simple circles via text)
for x_frac, y_frac, sz in [(0.05,0.1,8),(0.9,0.15,6),(0.15,0.85,5),
                             (0.8,0.8,7),(0.5,0.05,5),(0.35,0.92,6)]:
    star = s.shapes.add_shape(9, Inches(x_frac*13.33)-Pt(sz/2),
                               Inches(y_frac*7.5)-Pt(sz/2), Pt(sz), Pt(sz))
    star.fill.solid(); star.fill.fore_color.rgb = WHITE
    star.line.fill.background()

textbox(s, "🚀  Space Defender", Inches(1), Inches(1.6), Inches(11.3), Inches(1.5),
        size=54, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
textbox(s, "Build a Space Shooter with Python & Pygame",
        Inches(1), Inches(3.0), Inches(11.3), Inches(0.7),
        size=26, color=WHITE, align=PP_ALIGN.CENTER)
divider(s, Inches(3.9), CYAN, 10)
textbox(s, "A beginner-friendly game programming tutorial",
        Inches(1), Inches(4.1), Inches(11.3), Inches(0.5),
        size=18, color=GREY, align=PP_ALIGN.CENTER)
textbox(s, "10 slides  ·  ~300 lines of Python  ·  No image files needed",
        Inches(1), Inches(4.7), Inches(11.3), Inches(0.4),
        size=16, color=GREY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — What We're Building
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
textbox(s, "What We're Building", Inches(0.67), Inches(0.3), Inches(10), Inches(0.7),
        size=36, bold=True, color=CYAN)
divider(s, Inches(1.1))

# Left column — features
bullet_list(s, [
    "Player spaceship — move in all 4 directions, shoot laser bolts",
    "3 alien enemy types — each with different HP and points",
    "Enemies that dive at you and fire back",
    "Boss battle every 5th wave (30 HP, 3 cannons)",
    "Particle explosions on every kill",
    "Scrolling star-field background",
    "Lives system + invincibility flicker after being hit",
    "Synthesised sound effects (no audio files needed)",
], Inches(0.67), Inches(1.3), Inches(6), Inches(5.5), size=18)

# Right column — tech stack box
box = s.shapes.add_shape(1, Inches(7.2), Inches(1.3), Inches(5.5), Inches(5.5))
box.fill.solid(); box.fill.fore_color.rgb = DK_GREY
box.line.color.rgb = YELLOW; box.line.width = Pt(1.5)

textbox(s, "Tech Stack", Inches(7.4), Inches(1.5), Inches(5), Inches(0.5),
        size=18, bold=True, color=YELLOW)
bullet_list(s, [
    "Python 3",
    "pygame  (graphics, input, sprites)",
    "numpy   (sound synthesis)",
    "random / math  (game logic)",
    "Single file — space_game.py",
], Inches(7.4), Inches(2.1), Inches(5), Inches(4.5), size=17, indent_color=YELLOW)

textbox(s, "Install:", Inches(7.4), Inches(5.6), Inches(2), Inches(0.4),
        size=15, color=GREY)
code_box(s, "pip install pygame numpy", Inches(7.4), Inches(6.0), Inches(5.2), Inches(0.7), size=14)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Project Structure & Setup
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
textbox(s, "Project Setup & Imports", Inches(0.67), Inches(0.3), Inches(10), Inches(0.7),
        size=36, bold=True, color=CYAN)
divider(s, Inches(1.1))

textbox(s, "Everything lives in one file — no folders, no assets.",
        Inches(0.67), Inches(1.2), Inches(12), Inches(0.4), size=18, color=GREY)

code_box(s, """\
import pygame          # game window, drawing, input, sprites
import random          # enemy behaviour, particle angles
import math            # sine-wave movement, dive angles
import sys             # clean exit
import numpy as np     # synthesise sound waves

pygame.init()
pygame.mixer.init(frequency=44100, size=-16, channels=2, buffer=512)

WIDTH, HEIGHT = 900, 700
FPS = 60

screen = pygame.display.set_mode((WIDTH, HEIGHT))
pygame.display.set_caption("Space Defender")
clock = pygame.time.Clock()

# Colour palette — used everywhere
CYAN  = (0, 255, 255)
RED   = (220, 40, 40)
DARK  = (10, 10, 30)    # background colour""",
Inches(0.67), Inches(1.7), Inches(12), Inches(5.3), size=14)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — The Game Loop
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
textbox(s, "The Game Loop", Inches(0.67), Inches(0.3), Inches(10), Inches(0.7),
        size=36, bold=True, color=CYAN)
divider(s, Inches(1.1))

textbox(s, "Every game runs the same 4-step loop, 60 times per second:",
        Inches(0.67), Inches(1.2), Inches(12), Inches(0.4), size=18, color=WHITE)

# Loop diagram boxes
steps = [
    ("1  Handle Events", "Key presses, mouse clicks,\nwindow close", CYAN),
    ("2  Update", "Move sprites, check\ncollisions, apply physics", YELLOW),
    ("3  Draw", "Clear screen, draw all\nsprites and HUD", ORANGE),
    ("4  Flip", "Show the new frame\nclock.tick(60)", GREEN),
]
for i, (title, desc, col) in enumerate(steps):
    bx = Inches(0.5 + i * 3.2)
    box = s.shapes.add_shape(1, bx, Inches(1.75), Inches(3.0), Inches(1.9))
    box.fill.solid(); box.fill.fore_color.rgb = DK_GREY
    box.line.color.rgb = col; box.line.width = Pt(2)
    textbox(s, title, bx + Inches(0.1), Inches(1.85), Inches(2.8), Inches(0.5),
            size=16, bold=True, color=col)
    textbox(s, desc,  bx + Inches(0.1), Inches(2.35), Inches(2.8), Inches(1.2),
            size=14, color=WHITE)
    if i < 3:
        textbox(s, "→", Inches(3.35 + i * 3.2), Inches(2.3), Inches(0.5), Inches(0.5),
                size=28, color=GREY, align=PP_ALIGN.CENTER)

code_box(s, """\
running = True
while running:
    clock.tick(FPS)                        # ← cap at 60 FPS

    for event in pygame.event.get():       # 1. events
        if event.type == pygame.QUIT:
            running = False

    keys = pygame.key.get_pressed()
    player.update(keys)                    # 2. update
    all_sprites.update()

    screen.fill(DARK)                      # 3. draw
    all_sprites.draw(screen)

    pygame.display.flip()                  # 4. flip""",
Inches(0.67), Inches(3.85), Inches(12), Inches(3.3), size=13)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Sprites: Drawing Without Image Files
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
textbox(s, "Drawing Sprites with Code", Inches(0.67), Inches(0.3), Inches(10), Inches(0.7),
        size=36, bold=True, color=CYAN)
divider(s, Inches(1.1))

textbox(s, "Instead of loading PNG files, we draw shapes directly onto a Surface:",
        Inches(0.67), Inches(1.2), Inches(12), Inches(0.4), size=17, color=WHITE)

code_box(s, """\
def draw_player_ship(surface, color=CYAN):
    w, h = surface.get_size()
    cx, cy = w // 2, h // 2

    # Main body — a polygon (4 points)
    body = [(cx,       cy - h//2 + 4),   # nose
            (cx + w//3, cy + h//4),       # right wing tip
            (cx,       cy + h//6),        # engine centre
            (cx - w//3, cy + h//4)]       # left wing tip
    pygame.draw.polygon(surface, color, body)

    # Engine flame
    pygame.draw.polygon(surface, ORANGE,
        [(cx-8, cy+h//6), (cx+8, cy+h//6), (cx, cy+h//2-2)])

    # Cockpit bubble
    pygame.draw.ellipse(surface, WHITE, (cx-6, cy-14, 12, 14))

def make_surface(w, h, draw_fn, **kw):
    surf = pygame.Surface((w, h), pygame.SRCALPHA)  # transparent bg
    draw_fn(surf, **kw)
    return surf

# Usage:
player_image = make_surface(60, 70, draw_player_ship, color=CYAN)""",
Inches(0.67), Inches(1.7), Inches(12), Inches(5.2), size=13)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — The Player Class
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
textbox(s, "The Player Class", Inches(0.67), Inches(0.3), Inches(10), Inches(0.7),
        size=36, bold=True, color=CYAN)
divider(s, Inches(1.1))

# Left — code
code_box(s, """\
class Player(pygame.sprite.Sprite):
    def __init__(self):
        super().__init__()
        self.image = make_surface(60, 70, draw_player_ship)
        self.rect  = self.image.get_rect(
                         center=(WIDTH//2, HEIGHT-70))
        self.speed  = 6
        self.lives  = 3
        self.shield = 0   # invincibility frames
        self.shoot_cooldown = 0

    def update(self, keys):
        if keys[pygame.K_LEFT]:  self.rect.x -= self.speed
        if keys[pygame.K_RIGHT]: self.rect.x += self.speed
        if keys[pygame.K_UP]:    self.rect.y -= self.speed
        if keys[pygame.K_DOWN]:  self.rect.y += self.speed
        self.rect.clamp_ip(screen.get_rect()) # stay on screen

    def shoot(self):
        if self.shoot_cooldown == 0:
            self.shoot_cooldown = 15  # wait 15 frames
            return Bullet(self.rect.centerx,
                          self.rect.top + 10)
        return None

    def hit(self):
        if self.shield == 0:
            self.lives -= 1
            self.shield = 90  # ~1.5 sec invincible
            return True
        return False""",
Inches(0.67), Inches(1.2), Inches(7.5), Inches(5.8), size=12.5)

# Right — key concepts
textbox(s, "Key Concepts", Inches(8.5), Inches(1.3), Inches(4.5), Inches(0.45),
        size=18, bold=True, color=YELLOW)
bullet_list(s, [
    "Inherits pygame.sprite.Sprite",
    "self.image — what gets drawn",
    "self.rect  — position & hitbox",
    "clamp_ip() keeps ship in bounds",
    "Shoot cooldown prevents spam",
    "Shield = invincibility timer",
    "Controls: Arrow keys or WASD",
    "SPACE to fire",
], Inches(8.5), Inches(1.85), Inches(4.5), Inches(5.0), size=16, indent_color=YELLOW)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Enemies & Movement Patterns
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
textbox(s, "Enemies & Movement Patterns", Inches(0.67), Inches(0.3), Inches(12), Inches(0.7),
        size=36, bold=True, color=CYAN)
divider(s, Inches(1.1))

# Three enemy type boxes
types = [
    ("Saucer", "RED", "1 HP · 10 pts\nSlow, classic saucer shape", RED),
    ("Crab",   "ORANGE", "2 HP · 20 pts\nMedium speed, claw design", ORANGE),
    ("Fighter","PURPLE", "3 HP · 30 pts\nAngular, toughest enemy", PURPLE),
]
for i, (name, col_name, desc, col) in enumerate(types):
    bx = Inches(0.5 + i * 4.2)
    box = s.shapes.add_shape(1, bx, Inches(1.2), Inches(3.9), Inches(1.5))
    box.fill.solid(); box.fill.fore_color.rgb = DK_GREY
    box.line.color.rgb = col; box.line.width = Pt(2)
    textbox(s, name, bx + Inches(0.1), Inches(1.3), Inches(3.7), Inches(0.45),
            size=18, bold=True, color=col)
    textbox(s, desc, bx + Inches(0.1), Inches(1.75), Inches(3.7), Inches(0.8),
            size=14, color=WHITE)

code_box(s, """\
class Enemy(pygame.sprite.Sprite):
    def update(self):
        if not self.dive:
            # Formation: sine-wave left & right
            self.t += 0.025
            self.rect.x = int(self.origin_x + math.sin(self.t) * 30)
            self.rect.y += 0.3            # slowly drift down
        else:
            # Dive attack: fly straight toward player
            self.rect.x += int(self.dive_dx)
            self.rect.y += int(self.dive_dy)

    def start_dive(self, target_x):
        self.dive = True
        dx = target_x - self.rect.centerx
        dy = HEIGHT  - self.rect.y
        dist = math.hypot(dx, dy) or 1
        speed = self.base_speed * 4
        self.dive_dx = dx / dist * speed   # normalised direction × speed
        self.dive_dy = dy / dist * speed""",
Inches(0.67), Inches(2.9), Inches(12), Inches(4.2), size=13)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Bullets & Collision Detection
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
textbox(s, "Bullets & Collision Detection", Inches(0.67), Inches(0.3), Inches(12), Inches(0.7),
        size=36, bold=True, color=CYAN)
divider(s, Inches(1.1))

code_box(s, """\
class Bullet(pygame.sprite.Sprite):
    def __init__(self, x, y, dy=-16, color=CYAN, enemy=False):
        super().__init__()
        self.image = pygame.Surface((6, 18), pygame.SRCALPHA)
        pygame.draw.rect(self.image, color, (2, 0, 2, 18))  # laser beam
        pygame.draw.rect(self.image, WHITE,  (2, 2, 2, 6))  # bright tip
        self.rect = self.image.get_rect(center=(x, y))
        self.dy = dy          # negative = moves UP (player bullet)

    def update(self):
        self.rect.y += self.dy
        if self.rect.bottom < 0 or self.rect.top > HEIGHT:
            self.kill()       # remove when off-screen""",
Inches(0.67), Inches(1.2), Inches(12), Inches(3.1), size=13)

textbox(s, "Collision Detection — pygame does the hard work:",
        Inches(0.67), Inches(4.45), Inches(12), Inches(0.4), size=17, bold=True, color=YELLOW)

code_box(s, """\
# groupcollide(groupA, groupB, kill_A, kill_B)
# Returns dict: {sprite_from_A: [sprites_from_B_that_hit_it]}

hits = pygame.sprite.groupcollide(enemies, bullets, False, True)
for enemy, _ in hits.items():
    if enemy.take_hit():          # decrement HP
        explode(...)              # spawn particles
        player.score += enemy.pts
        enemy.kill()

# spritecollide checks one sprite against a whole group
if pygame.sprite.spritecollide(player, enemy_bullets, True):
    player.hit()                  # lose a life""",
Inches(0.67), Inches(4.9), Inches(12), Inches(2.4), size=13)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Particles & Sound
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
textbox(s, "Particles & Sound", Inches(0.67), Inches(0.3), Inches(10), Inches(0.7),
        size=36, bold=True, color=CYAN)
divider(s, Inches(1.1))

# Left — particles
textbox(s, "Explosion Particles", Inches(0.67), Inches(1.2), Inches(6), Inches(0.4),
        size=18, bold=True, color=ORANGE)
code_box(s, """\
class Particle(pygame.sprite.Sprite):
    def __init__(self, x, y, color):
        angle = random.uniform(0, 2 * math.pi)
        speed = random.uniform(1, 5)
        self.vx = math.cos(angle) * speed
        self.vy = math.sin(angle) * speed
        self.life = random.randint(20, 40)  # frames

    def update(self):
        self.life -= 1
        if self.life <= 0:
            self.kill()
        self.x += self.vx
        self.vy += 0.1   # gravity pulls down
        # Fade out as life decreases
        alpha = int(255 * self.life / 40)

def explode(groups, x, y, colors, count=30):
    for _ in range(count):
        p = Particle(x, y, random.choice(colors))
        for g in groups:
            g.add(p)""",
Inches(0.67), Inches(1.7), Inches(6.0), Inches(5.0), size=12.5)

# Right — sound
textbox(s, "Sound Synthesis with numpy", Inches(7.0), Inches(1.2), Inches(6), Inches(0.4),
        size=18, bold=True, color=GREEN)
code_box(s, """\
# No audio files — generate sounds with maths!

def make_laser_sound():
    t = np.linspace(0, 0.15,
            int(44100 * 0.15))
    # Frequency sweeps 900 → 300 Hz
    freq = np.linspace(900, 300, len(t))
    s = np.sin(2*np.pi
               * np.cumsum(freq)/44100)
    s *= np.linspace(1, 0, len(s))  # fade
    return pygame_sound(s)

# Each sound is just a numpy array:
# laser    = descending sine sweep
# explosion= noise × decay envelope
# hit      = low rumble + noise burst""",
Inches(7.0), Inches(1.7), Inches(6.0), Inches(3.7), size=12.5)

bullet_list(s, [
    "play('laser', 0.6)   — when player shoots",
    "play('explosion')    — on enemy death",
    "play('player_hit')   — when ship is hit",
    "play('wave_clear')   — ascending chime",
], Inches(7.0), Inches(5.55), Inches(6.0), Inches(1.8), size=15, indent_color=GREEN)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Putting It All Together
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
textbox(s, "Putting It All Together", Inches(0.67), Inches(0.3), Inches(12), Inches(0.7),
        size=36, bold=True, color=CYAN)
divider(s, Inches(1.1))

# State machine diagram
states = [
    ("TITLE",      CYAN,   Inches(0.5),  Inches(1.3)),
    ("PLAYING",    GREEN,  Inches(3.5),  Inches(1.3)),
    ("WAVE CLEAR", YELLOW, Inches(6.5),  Inches(1.3)),
    ("BOSS",       RED,    Inches(6.5),  Inches(2.6)),
    ("GAME OVER",  ORANGE, Inches(9.5),  Inches(1.3)),
]
for name, col, bx, by in states:
    box = s.shapes.add_shape(1, bx, by, Inches(2.7), Inches(0.9))
    box.fill.solid(); box.fill.fore_color.rgb = DK_GREY
    box.line.color.rgb = col; box.line.width = Pt(2)
    textbox(s, name, bx+Inches(0.05), by+Inches(0.2), Inches(2.6), Inches(0.5),
            size=16, bold=True, color=col, align=PP_ALIGN.CENTER)

# Arrows (as text — simple)
for text, l, t in [("ENTER →", Inches(3.15), Inches(1.65)),
                   ("clear →", Inches(6.2),  Inches(1.65)),
                   ("wave%5==0", Inches(7.8), Inches(2.25)),
                   ("no lives →", Inches(9.2), Inches(1.65)),
                   ("ENTER →", Inches(3.15), Inches(3.0))]:
    textbox(s, text, l, t, Inches(2.2), Inches(0.4), size=13, color=GREY)

# Summary code
code_box(s, """\
# Wave progression
def spawn_wave(wave, all_sprites, enemies):
    rows = min(2 + wave, 5)          # more rows each wave
    cols = min(5 + wave, 11)         # more columns each wave
    for row in range(rows):
        for col in range(cols):
            kind = min(row, 2)       # tougher enemies in lower rows
            e = Enemy(x, y, kind=kind)
            e.base_speed = 1.0 + wave * 0.15   # faster each wave
            enemies.add(e)

# Every 5th wave → Boss fight (30 HP, 3 cannons, spread shot)
if wave % 5 == 0:
    boss = Boss()""",
Inches(0.67), Inches(3.75), Inches(7.5), Inches(3.0), size=13)

# Recap bullets
textbox(s, "What You Learned", Inches(8.5), Inches(3.75), Inches(4.5), Inches(0.4),
        size=17, bold=True, color=YELLOW)
bullet_list(s, [
    "pygame Sprite class pattern",
    "The 4-step game loop",
    "Drawing shapes as sprites",
    "Collision detection groups",
    "Particle systems",
    "Procedural sound synthesis",
    "State machine game flow",
    "Wave-based difficulty scaling",
], Inches(8.5), Inches(4.25), Inches(4.5), Inches(2.9), size=15, indent_color=YELLOW)

textbox(s, "python3 space_game.py",
        Inches(0.67), Inches(6.9), Inches(12), Inches(0.45),
        size=18, bold=True, color=CYAN, align=PP_ALIGN.CENTER)

# ── Save ───────────────────────────────────────────────────────────────────────
out = "/Users/adil/Desktop/claude_test/space_game_slides.pptx"
prs.save(out)
print(f"Saved → {out}")
